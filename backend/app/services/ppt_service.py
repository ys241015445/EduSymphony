"""PPT rendering service using python-pptx with rich layouts and template metadata.

Supported layouts (16):
    title_slide | section_header | agenda | content |
    two_column | comparison | timeline | process_steps |
    quote | callout | stats | closing |
    big_number | quadrant | checklist | definition

Template metadata shape (all optional, resolved with sensible defaults):
    {
        "name":        "xxx",
        "mood":        "xxx",
        "palette":     {bg, title_color, body_color, accent, section_bg, bullet_color},
        "layout_style":"academic|modern|kawaii|tech|editorial|minimal|business|natural",
        "typography":  "serif|sans_display|handwriting|mono",
        "cover_style": "centered|split|decorative",
    }
"""
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Colour palettes per style (fallback if no AI palette given) ─────────

STYLES = {
    "academic": {
        "bg": RGBColor(0x1B, 0x2A, 0x4A),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "body_color": RGBColor(0xE0, 0xE0, 0xE0),
        "accent": RGBColor(0x4F, 0xC3, 0xF7),
        "section_bg": RGBColor(0x15, 0x22, 0x3B),
        "bullet_color": RGBColor(0x4F, 0xC3, 0xF7),
    },
    "modern": {
        "bg": RGBColor(0x0F, 0x0F, 0x23),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "body_color": RGBColor(0xCC, 0xCC, 0xCC),
        "accent": RGBColor(0x6C, 0x63, 0xFF),
        "section_bg": RGBColor(0x1A, 0x1A, 0x2E),
        "bullet_color": RGBColor(0x6C, 0x63, 0xFF),
    },
    "minimal": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x22, 0x22, 0x22),
        "body_color": RGBColor(0x44, 0x44, 0x44),
        "accent": RGBColor(0x00, 0x7A, 0xFF),
        "section_bg": RGBColor(0xF5, 0xF5, 0xF5),
        "bullet_color": RGBColor(0x00, 0x7A, 0xFF),
    },
    "colorful": {
        "bg": RGBColor(0xFF, 0xF8, 0xE1),
        "title_color": RGBColor(0xE6, 0x5C, 0x00),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0xFF, 0x6F, 0x00),
        "section_bg": RGBColor(0xFF, 0xEC, 0xB3),
        "bullet_color": RGBColor(0xE6, 0x5C, 0x00),
    },
}


TYPOGRAPHY_FONTS = {
    "serif":        {"title": "Source Han Serif SC",  "body": "Source Han Serif SC"},
    "sans_display": {"title": "Microsoft YaHei UI",   "body": "Microsoft YaHei"},
    "handwriting":  {"title": "楷体",                 "body": "楷体"},
    "mono":         {"title": "Consolas",             "body": "Microsoft YaHei"},
}

DEFAULT_FONTS = TYPOGRAPHY_FONTS["sans_display"]
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
PALETTE_KEYS = ("bg", "title_color", "body_color", "accent", "section_bg", "bullet_color")


# ── helpers ──────────────────────────────────────────────────────────────

def _hex_to_rgb(hex_str: str) -> RGBColor:
    try:
        s = (hex_str or "").strip().lstrip("#")
        if len(s) != 6:
            raise ValueError(f"bad hex length: {hex_str!r}")
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        return RGBColor(0xFF, 0xFF, 0xFF)


def _resolve_palette(style: str = "modern", palette: Optional[dict] = None) -> dict:
    if palette:
        return {k: _hex_to_rgb(palette.get(k, "")) for k in PALETTE_KEYS}
    return STYLES.get(style, STYLES["modern"])


def _resolve_fonts(template: Optional[dict]) -> dict:
    if not template:
        return DEFAULT_FONTS
    typ = (template.get("typography") or "").strip().lower()
    return TYPOGRAPHY_FONTS.get(typ, DEFAULT_FONTS)


def _cover_style(template: Optional[dict]) -> str:
    if not template:
        return "centered"
    cs = (template.get("cover_style") or "").strip().lower()
    return cs if cs in {"centered", "split", "decorative"} else "centered"


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size, color,
                  bold=False, alignment=PP_ALIGN.LEFT, font_name: Optional[str] = None):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or ""
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or DEFAULT_FONTS["body"]
    p.alignment = alignment
    return tx


def _add_rect(slide, left, top, width, height, fill_color: Optional[RGBColor] = None,
              line_color: Optional[RGBColor] = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_line(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_circle(slide, left, top, diameter, fill_color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _attach_notes(slide, notes: str):
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _title_bar(slide, title: str, pal: dict, fonts: dict):
    """Standard top-left page title + accent underline, shared by most layouts."""
    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
                  title or "", 30, pal["title_color"], bold=True, font_name=fonts["title"])
    _add_line(slide, Inches(0.8), Inches(1.3), Inches(1.8), Pt(3), pal["accent"])


def _bullet_list(slide, bullets, left, top, width, height, pal, fonts,
                 font_size=18, bullet_symbol="●"):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    items = [b for b in (bullets or []) if b]
    for i, b in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{bullet_symbol}  {b}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = pal["body_color"]
        p.font.name = fonts["body"]
        p.space_after = Pt(8)


# ── Renderers: existing four (with cover variants + typography) ─────────

def _render_title_slide(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    cs = _cover_style(tpl)
    title = sd.get("title", "") or ""
    subtitle = sd.get("subtitle", "") or ""

    if cs == "split":
        # Left colour block + right stacked type.
        _add_rect(slide, 0, 0, Inches(5.2), SLIDE_HEIGHT, fill_color=pal["section_bg"])
        _add_line(slide, Inches(0.7), Inches(2.8), Inches(3), Pt(6), pal["accent"])
        _add_text_box(slide, Inches(5.7), Inches(2.3), Inches(7), Inches(1.8),
                      title, 44, pal["title_color"], bold=True, font_name=fonts["title"])
        if subtitle:
            _add_text_box(slide, Inches(5.7), Inches(4.1), Inches(7), Inches(1.0),
                          subtitle, 20, pal["body_color"], font_name=fonts["body"])
    elif cs == "decorative":
        # Decorative circles + centered heading.
        _add_circle(slide, Inches(0.3), Inches(0.3), Inches(1.2), pal["accent"])
        _add_circle(slide, Inches(11.8), Inches(5.8), Inches(1.4), pal["accent"])
        _add_circle(slide, Inches(11.3), Inches(0.6), Inches(0.6), pal["bullet_color"])
        _add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.8),
                      title, 46, pal["title_color"], bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
        if subtitle:
            _add_text_box(slide, Inches(2), Inches(4.5), Inches(9), Inches(1),
                          subtitle, 22, pal["body_color"],
                          alignment=PP_ALIGN.CENTER, font_name=fonts["body"])
    else:  # centered (default)
        _add_text_box(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(1.5),
                      title, 46, pal["title_color"], bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
        _add_line(slide, Inches(5.2), Inches(4.0), Inches(3), Pt(4), pal["accent"])
        if subtitle:
            _add_text_box(slide, Inches(2), Inches(4.4), Inches(9), Inches(1),
                          subtitle, 22, pal["body_color"],
                          alignment=PP_ALIGN.CENTER, font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


def _render_section_header(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["section_bg"])
    fonts = _resolve_fonts(tpl)
    title = sd.get("title", "") or ""
    subtitle = sd.get("subtitle", "") or ""

    _add_line(slide, Inches(1.5), Inches(2.3), Inches(3), Pt(6), pal["accent"])
    _add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.3),
                  title, 40, pal["accent"], bold=True, font_name=fonts["title"])
    if subtitle:
        _add_text_box(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(1),
                      subtitle, 22, pal["body_color"], font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


def _render_content_slide(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)
    _bullet_list(slide, sd.get("bullets", []), Inches(1.0), Inches(1.8),
                 Inches(11), Inches(5), pal, fonts, font_size=18)
    _attach_notes(slide, sd.get("notes", ""))


def _render_closing_slide(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    title = sd.get("title", "谢谢") or "谢谢"
    _add_text_box(slide, Inches(2), Inches(2.6), Inches(9), Inches(2),
                  title, 52, pal["accent"], bold=True,
                  alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
    bullets = sd.get("bullets", []) or []
    if bullets:
        _add_text_box(slide, Inches(2), Inches(4.6), Inches(9), Inches(1.5),
                      "\n".join(bullets), 18, pal["body_color"],
                      alignment=PP_ALIGN.CENTER, font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


# ── Renderers: 8 new layouts ────────────────────────────────────────────

def _render_agenda(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", "目录"), pal, fonts)

    items = [b for b in (sd.get("bullets") or []) if b]
    top = Inches(1.9)
    for i, item in enumerate(items[:7]):
        y = top + Inches(i * 0.7)
        _add_text_box(slide, Inches(1.0), y, Inches(0.8), Inches(0.6),
                      f"{i + 1:02d}", 28, pal["accent"], bold=True,
                      font_name=fonts["title"])
        _add_text_box(slide, Inches(2.0), y + Inches(0.05), Inches(10), Inches(0.6),
                      str(item), 20, pal["body_color"], font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


def _render_two_column(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    left_items = sd.get("left_bullets") or []
    right_items = sd.get("right_bullets") or []
    if not left_items and not right_items:
        # fallback: split bullets in half
        bullets = [b for b in (sd.get("bullets") or []) if b]
        mid = (len(bullets) + 1) // 2
        left_items, right_items = bullets[:mid], bullets[mid:]

    _bullet_list(slide, left_items, Inches(0.9), Inches(1.9),
                 Inches(5.8), Inches(5), pal, fonts, font_size=17)
    _bullet_list(slide, right_items, Inches(6.9), Inches(1.9),
                 Inches(5.8), Inches(5), pal, fonts, font_size=17)
    _attach_notes(slide, sd.get("notes", ""))


def _render_comparison(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    left_title = sd.get("left_title", "") or "A"
    right_title = sd.get("right_title", "") or "B"
    left_items = sd.get("left_bullets") or []
    right_items = sd.get("right_bullets") or []

    card_top = Inches(1.9)
    card_h = Inches(5.1)
    # Left card
    _add_rect(slide, Inches(0.7), card_top, Inches(5.9), card_h,
              fill_color=pal["section_bg"], line_color=pal["accent"])
    _add_text_box(slide, Inches(1.0), Inches(2.0), Inches(5.4), Inches(0.7),
                  left_title, 22, pal["accent"], bold=True, font_name=fonts["title"])
    _bullet_list(slide, left_items, Inches(1.0), Inches(2.8),
                 Inches(5.4), Inches(4), pal, fonts, font_size=16, bullet_symbol="▸")
    # Right card
    _add_rect(slide, Inches(6.7), card_top, Inches(5.9), card_h,
              fill_color=pal["section_bg"], line_color=pal["bullet_color"])
    _add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.4), Inches(0.7),
                  right_title, 22, pal["bullet_color"], bold=True, font_name=fonts["title"])
    _bullet_list(slide, right_items, Inches(7.0), Inches(2.8),
                 Inches(5.4), Inches(4), pal, fonts, font_size=16, bullet_symbol="▸")
    _attach_notes(slide, sd.get("notes", ""))


def _render_timeline(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    steps = [s for s in (sd.get("steps") or []) if isinstance(s, dict)]
    n = min(len(steps), 6) or 1
    # horizontal line
    track_y = Inches(3.6)
    _add_line(slide, Inches(1.0), track_y + Inches(0.35), Inches(11.3), Pt(2), pal["accent"])
    seg_w = 11.3 / max(n, 1)
    for i, st in enumerate(steps[:n]):
        cx = 1.0 + seg_w * (i + 0.5) - 0.3
        _add_circle(slide, Inches(cx), track_y, Inches(0.6), pal["accent"])
        _add_text_box(slide, Inches(cx - 0.3), track_y + Inches(0.08), Inches(1.2), Inches(0.5),
                      str(i + 1), 18, pal["bg"], bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
        _add_text_box(slide, Inches(cx - 1.0), track_y - Inches(0.9), Inches(2.6), Inches(0.6),
                      str(st.get("name", ""))[:24], 16, pal["title_color"], bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
        desc = str(st.get("desc", ""))[:60]
        if desc:
            _add_text_box(slide, Inches(cx - 1.2), track_y + Inches(1.0), Inches(3.0), Inches(1.5),
                          desc, 12, pal["body_color"],
                          alignment=PP_ALIGN.CENTER, font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


def _render_process_steps(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    steps = [s for s in (sd.get("steps") or []) if isinstance(s, dict)][:6]
    top = Inches(1.9)
    for i, st in enumerate(steps):
        y = top + Inches(i * 0.85)
        _add_circle(slide, Inches(0.9), y + Inches(0.05), Inches(0.6), pal["accent"])
        _add_text_box(slide, Inches(0.9), y + Inches(0.12), Inches(0.6), Inches(0.5),
                      str(i + 1), 18, pal["bg"], bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
        _add_text_box(slide, Inches(1.7), y, Inches(4), Inches(0.7),
                      str(st.get("name", ""))[:30], 20, pal["title_color"], bold=True,
                      font_name=fonts["title"])
        _add_text_box(slide, Inches(5.7), y + Inches(0.1), Inches(7), Inches(0.7),
                      str(st.get("desc", ""))[:120], 15, pal["body_color"],
                      font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


def _render_quote(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["section_bg"])
    fonts = _resolve_fonts(tpl)

    quote_text = sd.get("quote") or sd.get("title") or ""
    author = sd.get("quote_author") or sd.get("subtitle") or ""
    _add_text_box(slide, Inches(1.2), Inches(1.5), Inches(1.5), Inches(1.5),
                  "“", 96, pal["accent"], bold=True, font_name=fonts["title"])
    _add_text_box(slide, Inches(1.8), Inches(2.5), Inches(9.7), Inches(3),
                  quote_text, 26, pal["title_color"], bold=True,
                  alignment=PP_ALIGN.LEFT, font_name=fonts["title"])
    if author:
        _add_line(slide, Inches(1.8), Inches(5.8), Inches(0.6), Pt(2), pal["accent"])
        _add_text_box(slide, Inches(2.5), Inches(5.6), Inches(9), Inches(0.6),
                      f"— {author}", 16, pal["body_color"], font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


def _render_callout(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    callout = sd.get("callout") or sd.get("subtitle") or ""
    _add_rect(slide, Inches(1.2), Inches(2.1), Inches(10.9), Inches(2.0),
              fill_color=pal["section_bg"], line_color=pal["accent"])
    _add_line(slide, Inches(1.2), Inches(2.1), Pt(8), Inches(2.0), pal["accent"])
    _add_text_box(slide, Inches(1.6), Inches(2.4), Inches(10.3), Inches(1.6),
                  callout, 22, pal["title_color"], bold=True, font_name=fonts["title"])

    bullets = [b for b in (sd.get("bullets") or []) if b]
    if bullets:
        _bullet_list(slide, bullets, Inches(1.2), Inches(4.5),
                     Inches(11), Inches(2.7), pal, fonts, font_size=16)
    _attach_notes(slide, sd.get("notes", ""))


def _render_stats(prs, sd, pal, tpl):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    stats = [s for s in (sd.get("stats") or []) if isinstance(s, dict)][:4]
    n = max(len(stats), 1)
    col_w = 11.3 / n
    top = Inches(2.4)
    for i, st in enumerate(stats):
        x = 1.0 + col_w * i
        _add_rect(slide, Inches(x), top, Inches(col_w - 0.3), Inches(3.4),
                  fill_color=pal["section_bg"])
        _add_text_box(slide, Inches(x), top + Inches(0.4), Inches(col_w - 0.3), Inches(1.4),
                      str(st.get("value", ""))[:10], 54, pal["accent"], bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
        _add_text_box(slide, Inches(x), top + Inches(2.0), Inches(col_w - 0.3), Inches(1.2),
                      str(st.get("label", ""))[:60], 15, pal["body_color"],
                      alignment=PP_ALIGN.CENTER, font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


# ── Renderers: 4 new layouts distilled from open-source PPT skills ───────
# big_number / quadrant / checklist / definition —— 强化数据、分类、清单、概念表达。

def _render_big_number(prs, sd, pal, tpl):
    """One hero number + short interpretation + optional supporting bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    number = str(sd.get("big_number") or "").strip()
    if not number:
        stats = [s for s in (sd.get("stats") or []) if isinstance(s, dict)]
        if stats:
            number = str(stats[0].get("value", ""))
    label = sd.get("big_label") or sd.get("subtitle") or ""
    if not label:
        stats = [s for s in (sd.get("stats") or []) if isinstance(s, dict)]
        if stats:
            label = str(stats[0].get("label", ""))

    _add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.2),
                  number[:12], 130, pal["accent"], bold=True,
                  alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
    if label:
        _add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(1.0),
                      str(label)[:60], 24, pal["title_color"], bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
    bullets = [b for b in (sd.get("bullets") or []) if b]
    if bullets:
        _bullet_list(slide, bullets[:3], Inches(2.4), Inches(5.5),
                     Inches(8.5), Inches(1.7), pal, fonts, font_size=15)
    _attach_notes(slide, sd.get("notes", ""))


def _render_quadrant(prs, sd, pal, tpl):
    """2x2 matrix; reads `quadrants` (4 items) or falls back to steps/bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    quads = [q for q in (sd.get("quadrants") or []) if isinstance(q, dict)]
    if not quads:
        steps = [s for s in (sd.get("steps") or []) if isinstance(s, dict)]
        quads = steps[:4]
    if not quads:
        bullets = [b for b in (sd.get("bullets") or []) if b][:4]
        quads = [{"name": "", "desc": b} for b in bullets]
    quads = (quads + [{}, {}, {}, {}])[:4]

    positions = [
        (Inches(0.8), Inches(1.9)), (Inches(6.9), Inches(1.9)),
        (Inches(0.8), Inches(4.5)), (Inches(6.9), Inches(4.5)),
    ]
    card_w, card_h = Inches(5.6), Inches(2.4)
    accents = [pal["accent"], pal["bullet_color"], pal["bullet_color"], pal["accent"]]
    for (x, y), q, acc in zip(positions, quads, accents):
        _add_rect(slide, x, y, card_w, card_h,
                  fill_color=pal["section_bg"], line_color=acc)
        name = str(q.get("name", "") or "")[:24]
        if name:
            _add_text_box(slide, x + Inches(0.25), y + Inches(0.15), card_w - Inches(0.5),
                          Inches(0.6), name, 18, acc, bold=True, font_name=fonts["title"])
        desc = str(q.get("desc", "") or "")[:110]
        if desc:
            _add_text_box(slide, x + Inches(0.25), y + Inches(0.85), card_w - Inches(0.5),
                          Inches(1.4), desc, 14, pal["body_color"], font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


def _render_checklist(prs, sd, pal, tpl):
    """Checkmark list for review pages / rules / reminders."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    items = [b for b in (sd.get("bullets") or []) if b][:7]
    top = Inches(1.9)
    for i, item in enumerate(items):
        y = top + Inches(i * 0.75)
        box = _add_rect(slide, Inches(1.0), y, Inches(0.45), Inches(0.45),
                        fill_color=pal["accent"])
        _add_text_box(slide, Inches(1.0), y - Inches(0.02), Inches(0.45), Inches(0.5),
                      "✓", 20, pal["bg"], bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=fonts["title"])
        _add_text_box(slide, Inches(1.7), y - Inches(0.02), Inches(10.6), Inches(0.6),
                      str(item)[:80], 18, pal["body_color"], font_name=fonts["body"])
    _attach_notes(slide, sd.get("notes", ""))


def _render_definition(prs, sd, pal, tpl):
    """Term + precise definition emphasised, with optional expansion bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, pal["bg"])
    fonts = _resolve_fonts(tpl)
    _title_bar(slide, sd.get("title", ""), pal, fonts)

    term = sd.get("term") or sd.get("subtitle") or ""
    definition = sd.get("definition") or ""
    _add_rect(slide, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.4),
              fill_color=pal["section_bg"], line_color=pal["accent"])
    _add_line(slide, Inches(0.8), Inches(1.9), Pt(8), Inches(2.4), pal["accent"])
    if term:
        _add_text_box(slide, Inches(1.2), Inches(2.15), Inches(11), Inches(0.8),
                      str(term)[:40], 30, pal["accent"], bold=True, font_name=fonts["title"])
    if definition:
        _add_text_box(slide, Inches(1.2), Inches(3.0), Inches(11), Inches(1.2),
                      str(definition)[:160], 18, pal["title_color"], font_name=fonts["body"])
    bullets = [b for b in (sd.get("bullets") or []) if b]
    if bullets:
        _bullet_list(slide, bullets[:4], Inches(1.0), Inches(4.6),
                     Inches(11.3), Inches(2.6), pal, fonts, font_size=16)
    _attach_notes(slide, sd.get("notes", ""))


# ── Dispatch + public entrypoints ───────────────────────────────────────

_LAYOUT_RENDERERS = {
    "title_slide":   _render_title_slide,
    "section_header": _render_section_header,
    "agenda":        _render_agenda,
    "content":       _render_content_slide,
    "two_column":    _render_two_column,
    "comparison":    _render_comparison,
    "timeline":      _render_timeline,
    "process_steps": _render_process_steps,
    "quote":         _render_quote,
    "callout":       _render_callout,
    "stats":         _render_stats,
    "big_number":    _render_big_number,
    "quadrant":      _render_quadrant,
    "checklist":     _render_checklist,
    "definition":    _render_definition,
    "closing":       _render_closing_slide,
}


def _render_slide(prs, sd, pal, tpl):
    layout = (sd.get("layout") or "content").strip().lower()
    fn = _LAYOUT_RENDERERS.get(layout, _render_content_slide)
    try:
        fn(prs, sd, pal, tpl)
    except Exception:
        # Last-resort: render as plain content so the whole deck doesn't die on 1 bad slide.
        _render_content_slide(prs, sd, pal, tpl)


def build_pptx(data: dict, style: str = "modern",
               palette: Optional[dict] = None,
               template: Optional[dict] = None) -> bytes:
    pal = _resolve_palette(style, palette or (template.get("palette") if template else None))
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slides = data.get("slides", []) if isinstance(data, dict) else []
    if not slides:
        slides = [{
            "layout": "title_slide",
            "title": (data or {}).get("title", "演示文稿"),
            "subtitle": (data or {}).get("subtitle", ""),
            "bullets": [], "notes": "",
        }]

    for sd in slides:
        if isinstance(sd, dict):
            _render_slide(prs, sd, pal, template)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def append_slides_to_pptx(existing_path: str, new_slides: list,
                          style: str = "modern",
                          palette: Optional[dict] = None,
                          template: Optional[dict] = None) -> bytes:
    pal = _resolve_palette(style, palette or (template.get("palette") if template else None))
    prs = Presentation(existing_path)
    for sd in (new_slides or []):
        if isinstance(sd, dict):
            _render_slide(prs, sd, pal, template)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
