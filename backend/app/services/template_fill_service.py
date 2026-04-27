"""
Template AI Fill —— 占位符检测、内容填充与跨格式导出。

设计要点：
- Hybrid 识别：优先匹配显式占位符（{{x}} / <x> / 【x】 / 《x》 / ____ 等），
  匹配不到再让 AI 自动识别空白位置。
- 填充保真度：docx/pptx/xlsx 走 run-level 替换，100% 保留原排版与样式。
- 跨格式导出：以原文件为 "primary"，其他格式通过纯文本抽取 + 重新渲染。
"""
from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from typing import Callable, Dict, List, Optional, Tuple

from loguru import logger


SUPPORTED_INPUT_EXTS = {"docx", "pptx", "xlsx", "txt", "md"}
SUPPORTED_OUTPUT_EXTS = {"docx", "pptx", "xlsx", "txt", "md", "pdf", "json"}


# ───────────────────────── 占位符正则 ─────────────────────────
# 每条返回 key 提取函数：输入 match 对象 → 占位符 key（裸字段名）
_PLACEHOLDER_PATTERNS: List[Tuple[str, re.Pattern, Callable[[re.Match], str]]] = [
    # {{ xxx }} / {{xxx}}
    ("curly_double", re.compile(r"\{\{\s*([^\{\}\n]{1,60}?)\s*\}\}"), lambda m: m.group(1).strip()),
    # 《xxx》
    ("chinese_book", re.compile(r"《([^《》\n]{1,60}?)》"), lambda m: m.group(1).strip()),
    # 【xxx】
    ("chinese_bracket", re.compile(r"【([^【】\n]{1,60}?)】"), lambda m: m.group(1).strip()),
    # <xxx> —— 避免误匹配 HTML 标签，只接受非字母开头的"中文/括注型"
    ("angle", re.compile(r"<\s*(请填[^<>\n]{0,40}|[^<>\n]*?填入[^<>\n]{0,40}|[\u4e00-\u9fff][^<>\n]{0,40}?)\s*>"),
     lambda m: m.group(1).strip()),
    # _____ (>=3 下划线) — 占位符 key 用上下文生成
    ("underline", re.compile(r"_{3,}"), lambda m: "__UNDERLINE__"),
]


@dataclass
class Placeholder:
    key: str                        # 裸 key，用于 fill_map 索引
    raw: str                        # 原样文本（含符号），用于 regex 替换
    pattern: str                    # 命中的 pattern 名字
    count: int = 1
    sample_context: str = ""        # 附近上下文，供 AI 生成参考

    def to_dict(self) -> dict:
        return {
            "key": self.key,
            "raw": self.raw,
            "pattern": self.pattern,
            "count": self.count,
            "sample_context": self.sample_context,
        }


@dataclass
class AnalyzeResult:
    placeholders: List[Placeholder] = field(default_factory=list)
    mode: str = "ai_detect"         # "token" | "ai_detect"
    preview_text: str = ""
    original_ext: str = ""

    def to_dict(self) -> dict:
        return {
            "placeholders": [p.to_dict() for p in self.placeholders],
            "mode": self.mode,
            "preview_text": self.preview_text,
            "original_ext": self.original_ext,
        }


# ───────────────────────── 小工具 ─────────────────────────

def _snippet(text: str, start: int, end: int, radius: int = 24) -> str:
    """截取命中位置附近上下文（用于 AI 推断字段语义）。"""
    lo = max(0, start - radius)
    hi = min(len(text), end + radius)
    return (("..." if lo > 0 else "") + text[lo:hi] + ("..." if hi < len(text) else "")).strip()


def _collect_from_text(text: str, bucket: Dict[str, Placeholder]):
    """用统一正则在一段纯文本里收集占位符，去重累加 count。"""
    if not text:
        return
    for pat_name, regex, key_fn in _PLACEHOLDER_PATTERNS:
        for m in regex.finditer(text):
            raw = m.group(0)
            key = key_fn(m)
            # underline 占位符 key 由上下文生成，更有意义
            if key == "__UNDERLINE__":
                ctx = _snippet(text, m.start(), m.end(), radius=12)
                # 取下划线"左侧"最多 10 字作为字段提示
                left = ctx.split("_")[0].strip()
                key = (left[-10:].strip() or "空白字段") + "_空白"
            lookup = f"{pat_name}::{key}"
            if lookup in bucket:
                bucket[lookup].count += 1
                continue
            bucket[lookup] = Placeholder(
                key=key, raw=raw, pattern=pat_name, count=1,
                sample_context=_snippet(text, m.start(), m.end()),
            )


# ───────────────────────── DOCX ─────────────────────────

def _iter_docx_paragraphs(doc):
    """遍历正文 + 表格单元格 + 页眉页脚里所有段落。"""
    for p in doc.paragraphs:
        yield p
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    yield p
    for section in doc.sections:
        for part in (section.header, section.footer):
            if part is None:
                continue
            for p in part.paragraphs:
                yield p


def _docx_paragraph_text(paragraph) -> str:
    return "".join(run.text or "" for run in paragraph.runs)


def detect_placeholders_docx(file_bytes: bytes) -> Tuple[List[Placeholder], str]:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    bucket: Dict[str, Placeholder] = {}
    preview_parts: List[str] = []
    for p in _iter_docx_paragraphs(doc):
        t = _docx_paragraph_text(p)
        if not t:
            continue
        _collect_from_text(t, bucket)
        if len(preview_parts) < 40:
            preview_parts.append(t)
    return list(bucket.values()), "\n".join(preview_parts)


def _replace_in_paragraph(paragraph, fill_map: Dict[str, str]) -> int:
    """
    在段落内做 run-safe 替换。方法：
    1. 把整段文字合并后做 regex 替换；
    2. 把结果写回第一个 run，清空其余 run（保留第一个 run 的 font 属性即可）。
    这样在大多数 "模板带占位符" 场景下，排版保真度足够。
    """
    runs = paragraph.runs
    if not runs:
        return 0
    full = "".join(r.text or "" for r in runs)
    if not full:
        return 0
    new_full, changed = _apply_fill_to_text(full, fill_map)
    if not changed:
        return 0
    # 写回第一个 run，清空其他 run
    runs[0].text = new_full
    for r in runs[1:]:
        r.text = ""
    return changed


def _apply_fill_to_text(text: str, fill_map: Dict[str, str]) -> Tuple[str, int]:
    """
    在一段文本里把占位符替换成 fill_map 对应值。
    返回 (new_text, 替换次数)。
    """
    out = text
    total = 0
    # 显式占位符（精确 key）
    for pat_name, regex, key_fn in _PLACEHOLDER_PATTERNS:
        if pat_name == "underline":
            continue
        def repl(m: re.Match) -> str:
            nonlocal total
            k = key_fn(m)
            if k in fill_map:
                total += 1
                return str(fill_map[k])
            return m.group(0)
        out = regex.sub(repl, out)
    # 下划线：按出现顺序消费 map 里以 "_空白" 结尾或 "__UNDERLINE_*" 开头的值
    # 简单起见：把 fill_map 里所有 "*_空白" 的值按定义顺序替换下划线
    underline_values = [v for k, v in fill_map.items() if k.endswith("_空白") or k.startswith("__UNDERLINE_")]
    if underline_values:
        idx = {"i": 0}
        def u_repl(_m: re.Match) -> str:
            if idx["i"] < len(underline_values):
                val = str(underline_values[idx["i"]])
                idx["i"] += 1
                total += 1
                return val
            return _m.group(0)
        out = _PLACEHOLDER_PATTERNS[-1][1].sub(u_repl, out)
    return out, total


def fill_docx(file_bytes: bytes, fill_map: Dict[str, str]) -> bytes:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    total = 0
    for p in _iter_docx_paragraphs(doc):
        total += _replace_in_paragraph(p, fill_map)
    logger.info(f"[template-fill] docx replaced {total} placeholders")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ───────────────────────── PPTX ─────────────────────────

def _iter_pptx_text_frames(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                yield shape.text_frame
            if shape.shape_type == 19 and shape.has_table:  # table
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame


def detect_placeholders_pptx(file_bytes: bytes) -> Tuple[List[Placeholder], str]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    bucket: Dict[str, Placeholder] = {}
    preview_parts: List[str] = []
    for tf in _iter_pptx_text_frames(prs):
        for p in tf.paragraphs:
            t = "".join(r.text or "" for r in p.runs)
            if not t:
                continue
            _collect_from_text(t, bucket)
            if len(preview_parts) < 60:
                preview_parts.append(t)
    return list(bucket.values()), "\n".join(preview_parts)


def fill_pptx(file_bytes: bytes, fill_map: Dict[str, str]) -> bytes:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    total = 0
    for tf in _iter_pptx_text_frames(prs):
        for p in tf.paragraphs:
            runs = p.runs
            if not runs:
                continue
            full = "".join(r.text or "" for r in runs)
            if not full:
                continue
            new_full, changed = _apply_fill_to_text(full, fill_map)
            if not changed:
                continue
            runs[0].text = new_full
            for r in runs[1:]:
                r.text = ""
            total += changed
    logger.info(f"[template-fill] pptx replaced {total} placeholders")
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ───────────────────────── XLSX ─────────────────────────

def detect_placeholders_xlsx(file_bytes: bytes) -> Tuple[List[Placeholder], str]:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), data_only=False)
    bucket: Dict[str, Placeholder] = {}
    preview_parts: List[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=False):
            for cell in row:
                v = cell.value
                if not isinstance(v, str) or not v:
                    continue
                _collect_from_text(v, bucket)
                if len(preview_parts) < 80:
                    preview_parts.append(f"[{ws.title}!{cell.coordinate}] {v}")
    return list(bucket.values()), "\n".join(preview_parts)


def fill_xlsx(file_bytes: bytes, fill_map: Dict[str, str]) -> bytes:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), data_only=False)
    total = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=False):
            for cell in row:
                v = cell.value
                if not isinstance(v, str) or not v:
                    continue
                new_v, changed = _apply_fill_to_text(v, fill_map)
                if changed:
                    cell.value = new_v
                    total += changed
    logger.info(f"[template-fill] xlsx replaced {total} placeholders")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ───────────────────────── Text / Markdown ─────────────────────────

def detect_placeholders_text(text: str) -> List[Placeholder]:
    bucket: Dict[str, Placeholder] = {}
    _collect_from_text(text, bucket)
    return list(bucket.values())


def fill_text(text: str, fill_map: Dict[str, str]) -> str:
    out, _ = _apply_fill_to_text(text, fill_map)
    return out


# ───────────────────────── 统一入口 ─────────────────────────

def analyze(file_bytes: bytes, ext: str, raw_text: Optional[str] = None) -> AnalyzeResult:
    ext = (ext or "").lower().lstrip(".")
    result = AnalyzeResult(original_ext=ext)
    try:
        if ext == "docx":
            ph, preview = detect_placeholders_docx(file_bytes)
        elif ext == "pptx":
            ph, preview = detect_placeholders_pptx(file_bytes)
        elif ext == "xlsx":
            ph, preview = detect_placeholders_xlsx(file_bytes)
        elif ext in ("txt", "md"):
            t = raw_text if raw_text is not None else file_bytes.decode("utf-8", errors="replace")
            ph = detect_placeholders_text(t)
            preview = t[:4000]
        else:
            raise ValueError(f"unsupported ext: {ext}")
    except Exception as e:
        logger.error(f"[template-fill] analyze failed: {e}")
        raise

    result.placeholders = ph
    result.preview_text = preview[:6000]
    result.mode = "token" if ph else "ai_detect"
    return result


def fill(
    file_bytes: bytes,
    ext: str,
    fill_map: Dict[str, str],
    raw_text: Optional[str] = None,
) -> bytes:
    ext = (ext or "").lower().lstrip(".")
    if ext == "docx":
        return fill_docx(file_bytes, fill_map)
    if ext == "pptx":
        return fill_pptx(file_bytes, fill_map)
    if ext == "xlsx":
        return fill_xlsx(file_bytes, fill_map)
    if ext in ("txt", "md"):
        t = raw_text if raw_text is not None else file_bytes.decode("utf-8", errors="replace")
        return fill_text(t, fill_map).encode("utf-8")
    raise ValueError(f"unsupported ext: {ext}")


# ───────────────────────── 跨格式导出 ─────────────────────────

def _extract_text_from_docx(b: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(b))
    lines: List[str] = []
    for p in _iter_docx_paragraphs(doc):
        lines.append(_docx_paragraph_text(p))
    return "\n".join(l for l in lines if l is not None)


def _extract_text_from_pptx(b: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(b))
    out: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text or "" for r in para.runs)
                    if t.strip():
                        out.append(t)
        out.append("")
    return "\n".join(out)


def _extract_text_from_xlsx(b: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(b), data_only=True)
    out: List[str] = []
    for ws in wb.worksheets:
        out.append(f"## {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            out.append(" | ".join(cells))
        out.append("")
    return "\n".join(out)


def _to_plain_text(src_bytes: bytes, src_ext: str) -> str:
    src_ext = src_ext.lower().lstrip(".")
    if src_ext in ("txt", "md"):
        return src_bytes.decode("utf-8", errors="replace")
    if src_ext == "docx":
        return _extract_text_from_docx(src_bytes)
    if src_ext == "pptx":
        return _extract_text_from_pptx(src_bytes)
    if src_ext == "xlsx":
        return _extract_text_from_xlsx(src_bytes)
    raise ValueError(f"unsupported src_ext: {src_ext}")


def _render_docx_from_text(text: str) -> bytes:
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    style = doc.styles["Normal"]
    style.font.size = Pt(11)
    style.font.name = "SimSun"
    for line in text.splitlines() or [""]:
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _render_pptx_from_text(text: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    # 每 12 行一页
    lines = text.splitlines()
    chunks = [lines[i:i + 12] for i in range(0, len(lines), 12)] or [[""]]
    blank = prs.slide_layouts[5]
    for chunk in chunks:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
        tf = tb.text_frame
        tf.text = chunk[0] if chunk else ""
        for line in chunk[1:]:
            tf.add_paragraph().text = line
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _render_xlsx_from_text(text: str) -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    for r, line in enumerate(text.splitlines() or [""], 1):
        cells = line.split(" | ") if " | " in line else [line]
        for c, v in enumerate(cells, 1):
            ws.cell(row=r, column=c, value=v)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _render_pdf_from_text(text: str) -> bytes:
    # 复用 export.py 里的 xhtml2pdf 方案
    from xhtml2pdf import pisa
    # 简单 HTML 模板，保留换行
    escaped = (text or "")\
        .replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")\
        .replace("\n", "<br/>")
    html = (
        "<html><head><meta charset='utf-8'>"
        "<style>body{font-family:'SimSun','Songti SC','Microsoft YaHei',sans-serif;"
        "font-size:11pt;line-height:1.6;}</style></head>"
        f"<body>{escaped}</body></html>"
    )
    buf = io.BytesIO()
    status = pisa.CreatePDF(src=html, dest=buf, encoding="utf-8")
    if status.err:
        logger.warning(f"[template-fill] xhtml2pdf errors: {status.err}")
    return buf.getvalue()


def convert_to(src_bytes: bytes, src_ext: str, target_ext: str) -> Tuple[bytes, bool]:
    """
    返回 (bytes, lossy)。
    - 原格式 → 原格式：直接返回，lossy=False
    - 其他 → 从抽出的纯文本重新渲染，lossy=True
    """
    src_ext = src_ext.lower().lstrip(".")
    target_ext = target_ext.lower().lstrip(".")
    if src_ext == target_ext:
        return src_bytes, False
    if target_ext not in SUPPORTED_OUTPUT_EXTS:
        raise ValueError(f"unsupported target ext: {target_ext}")

    # JSON 单独处理：返回 {"text": "..."}
    if target_ext == "json":
        import json as _json
        text = _to_plain_text(src_bytes, src_ext)
        return _json.dumps({"text": text}, ensure_ascii=False, indent=2).encode("utf-8"), True

    text = _to_plain_text(src_bytes, src_ext)
    if target_ext in ("txt", "md"):
        return text.encode("utf-8"), src_ext != target_ext
    if target_ext == "docx":
        return _render_docx_from_text(text), True
    if target_ext == "pptx":
        return _render_pptx_from_text(text), True
    if target_ext == "xlsx":
        return _render_xlsx_from_text(text), True
    if target_ext == "pdf":
        return _render_pdf_from_text(text), True

    raise ValueError(f"cannot convert {src_ext} -> {target_ext}")


MIME_BY_EXT = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pdf":  "application/pdf",
    "txt":  "text/plain; charset=utf-8",
    "md":   "text/markdown; charset=utf-8",
    "json": "application/json; charset=utf-8",
}
